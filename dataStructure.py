"""
长文档（教材类）分块策略

核心思路：
1. 层级标题分块：识别章/节/小节标题，每个最小节作为一个语义单元
2. 上下文面包屑：每块元数据携带完整标题路径（章>节>小节），检索时可定位来源
3. 滑动窗口兜底：对超长小节用带重叠的段落感知窗口再切分，不在句中断开
4. 短块合并：相邻过短的块向下合并，避免碎片化
"""

import re
import json
import logging
import os
import tempfile
from pathlib import Path
from typing import List, Dict, Optional, Tuple

import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from tqdm import tqdm

from rag_config import DEVICE

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# ============================================================
# 配置
# ============================================================

CHUNK_SIZE = 800        # 目标块字符数（中文教材一段约200字，4段约800字）
CHUNK_OVERLAP = 150     # 滑动窗口重叠字符数
MIN_CHUNK_SIZE = 100    # 低于此值的块向下合并

# 中文教材标题模式，按优先级（级别从高到低）排列
# 每项：(正则, 标题级别)；级别数字越小层级越高
# level=-1 表示 Markdown 标题，特殊处理
HEADING_PATTERNS: List[Tuple[re.Pattern, int]] = [
    (re.compile(r'^第[一二三四五六七八九十百千万零]+编\s*(.*)$'), 0),   # 编（最高层）
    (re.compile(r'^第[一二三四五六七八九十百千万零]+章\s*(.*)$'), 1),   # 章
    (re.compile(r'^第[一二三四五六七八九十百千万零]+节\s*(.*)$'), 2),   # 节
    (re.compile(r'^[一二三四五六七八九十]+[、．.]\s*(.+)$'), 3),        # 一、二、
    (re.compile(r'^[（(][一二三四五六七八九十]+[）)]\s*(.+)$'), 4),     # （一）（二）
    (re.compile(r'^\d+[、．.]\s*(.+)$'), 5),                           # 1. 2.
    (re.compile(r'^(#{1,6})\s+(.+)$'), -1),                            # Markdown 标题
]


# ============================================================
# 文本提取
# ============================================================

def _extract_pdf_text(pdf_path: str) -> str:
    """优先用PyMuPDF直接提取文字层；若页面文字为空则用EasyOCR兜底。"""
    doc = fitz.open(pdf_path)
    pages_text = []
    ocr_needed = []

    for i, page in enumerate(doc):
        text = page.get_text("text").strip()
        if text:
            pages_text.append((i, text))
        else:
            ocr_needed.append(i)

    if ocr_needed:
        logger.info(f"{len(ocr_needed)} 页无文字层，使用EasyOCR")
        import easyocr
        reader = easyocr.Reader(["ch_sim", "en"], gpu=(DEVICE == "cuda"))
        matrix = fitz.Matrix(300 / 72, 300 / 72)
        for i in ocr_needed:
            page = doc[i]
            pix = page.get_pixmap(matrix=matrix, alpha=False)
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as f:
                tmp = f.name
            pix.save(tmp)
            try:
                result = reader.readtext(tmp)
                lines = [r[1].strip() for r in result if r[1].strip()]
                pages_text.append((i, "\n".join(lines)))
            finally:
                os.unlink(tmp)

    doc.close()
    pages_text.sort(key=lambda x: x[0])
    return "\n\n".join(t for _, t in pages_text)


def _extract_docx_text(docx_path: str) -> str:
    doc = Document(docx_path)
    lines = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = para.style.name
        if style.startswith('Heading'):
            try:
                level = int(style.split()[-1])
            except ValueError:
                level = 1
            lines.append('#' * level + ' ' + text)
        else:
            lines.append(text)
    return "\n".join(lines)


def _extract_pptx_text(pptx_path: str) -> str:
    prs = Presentation(pptx_path)
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        lines.append(t)
        lines.append('')
    return "\n".join(lines)


def extract_text(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()
    if ext == '.pdf':
        return _extract_pdf_text(file_path)
    elif ext == '.docx':
        return _extract_docx_text(file_path)
    elif ext == '.pptx':
        return _extract_pptx_text(file_path)
    elif ext == '.csv':
        # 每行一条，逗号分隔原样保留
        with open(file_path, encoding='utf-8') as f:
            return f.read()
    elif ext in ('.txt', '.md'):
        with open(file_path, encoding='utf-8') as f:
            return f.read()
    else:
        raise ValueError(f"不支持的文件格式: {ext}")


# ============================================================
# 标题识别
# ============================================================

def _match_heading(line: str) -> Optional[Tuple[int, str]]:
    """
    尝试将一行文本匹配为标题。
    返回 (level, title_text)，未匹配返回 None。
    """
    line = line.strip()
    if not line:
        return None
    for pattern, level in HEADING_PATTERNS:
        m = pattern.match(line)
        if m:
            if level == -1:  # Markdown 标题
                return (len(m.group(1)), m.group(2).strip())
            groups = m.groups()
            title = groups[-1].strip() if groups else line
            return (level, title or line)
    return None


def _preprocess(text: str) -> List[str]:
    """
    预处理：合并被 OCR 换行打断的标题行。
    例如 "第" + "一章 xxx" -> "第一章 xxx"
    """
    lines = text.splitlines()
    result = []
    i = 0
    while i < len(lines):
        stripped = lines[i].strip()
        if (len(stripped) <= 2
                and stripped
                and i + 1 < len(lines)
                and re.search(r'[第一二三四五六七八九十章节编]$', stripped)):
            merged = stripped + lines[i + 1].strip()
            if _match_heading(merged):
                result.append(merged)
                i += 2
                continue
        result.append(lines[i])
        i += 1
    return result


# ============================================================
# 超长节的二次切分（段落感知滑动窗口）
# ============================================================

def _split_long_section(
    content: str,
    chunk_size: int = CHUNK_SIZE,
    overlap: int = CHUNK_OVERLAP,
) -> List[str]:
    """
    对超长节内容进行二次切分。
    优先以空行（段落边界）为切分点；段落本身过长时再按句末标点切句。
    相邻块保留 overlap 字符重叠，维持上下文连贯性。
    """
    paragraphs = [p.strip() for p in re.split(r'\n{2}', content) if p.strip()]
    if not paragraphs:
        return [content.strip()] if content.strip() else []

    chunks: List[str] = []
    current: List[str] = []
    current_len = 0

    def flush(parts: List[str]) -> str:
        return '\n\n'.join(parts)

    for para in paragraphs:
        # 单段落本身超过 chunk_size：按句切分
        if len(para) > chunk_size:
            if current:
                chunks.append(flush(current))
                current, current_len = [], 0
            sentences = re.split(r'(?<=[。！？；.!?;])\s*', para)
            sentences = [s.strip() for s in sentences if s.strip()]
            buf, buf_len = [], 0
            for sent in sentences:
                if buf_len + len(sent) > chunk_size and buf:
                    chunk_text = ''.join(buf)
                    chunks.append(chunk_text)
                    tail = chunk_text[-overlap:]
                    buf, buf_len = [tail, sent], len(tail) + len(sent)
                else:
                    buf.append(sent)
                    buf_len += len(sent)
            if buf:
                chunks.append(''.join(buf))
            continue

        if current_len + len(para) > chunk_size and current:
            chunk_text = flush(current)
            chunks.append(chunk_text)
            # 段落粒度重叠：保留最后一段作为下一块的开头
            tail_para = current[-1]
            overlap_para = tail_para if len(tail_para) <= overlap * 2 else tail_para[-overlap:]
            current = [overlap_para, para]
            current_len = len(overlap_para) + len(para)
        else:
            current.append(para)
            current_len += len(para)

    if current:
        chunks.append(flush(current))

    return chunks


# ============================================================
# 短块合并
# ============================================================

def _merge_short_chunks(chunks: List[Dict], min_size: int = MIN_CHUNK_SIZE) -> List[Dict]:
    """将过短的块向后合并，最后一块过短则向前合并。"""
    if not chunks:
        return chunks
    merged: List[Dict] = []
    i = 0
    while i < len(chunks):
        chunk = chunks[i]
        if len(chunk['text']) < min_size and i + 1 < len(chunks):
            nxt = dict(chunks[i + 1])
            nxt['text'] = chunk['text'] + '\n\n' + nxt['text']
            chunks[i + 1] = nxt
            i += 1
            continue
        merged.append(chunk)
        i += 1
    if len(merged) >= 2 and len(merged[-1]['text']) < min_size:
        last = merged.pop()
        merged[-1] = dict(merged[-1])
        merged[-1]['text'] += '\n\n' + last['text']
    return merged


# ============================================================
# 主分块函数
# ============================================================

def chunk_textbook(
    text: str,
    source: str = '',
    chunk_size: int = CHUNK_SIZE,
    overlap: int = CHUNK_OVERLAP,
    min_chunk_size: int = MIN_CHUNK_SIZE,
) -> List[Dict]:
    """
    长文档（教材）分块主函数。

    分块策略：
    1. 逐行识别章/节/小节等层级标题，以最小标题节为分块单元。
    2. 每个块的元数据包含完整的标题路径（面包屑），便于检索时溯源。
    3. 超过 chunk_size 的节用带重叠的段落感知滑动窗口再切分。
    4. 过短的块（< min_chunk_size）向下合并，减少碎片。

    返回 List[Dict]，每项：
      - text:     块文本（包含标题行）
      - metadata: source / breadcrumb / heading / level / chunk_index / total_chunks
    """
    lines = _preprocess(text)

    # ---- 第一步：按标题行切分成节 ----
    sections: List[Dict] = []   # {level, heading, lines}
    # 层级栈：记录当前激活的各级标题，用于生成面包屑
    level_stack: Dict[int, str] = {}
    current: Optional[Dict] = None

    for line in lines:
        heading_info = _match_heading(line)
        if heading_info:
            level, title = heading_info
            # 保存当前节
            if current is not None:
                sections.append(current)
            # 更新层级栈：清除比当前级别低（数字更大）的层级
            level_stack = {k: v for k, v in level_stack.items() if k < level}
            level_stack[level] = title
            # 生成面包屑：按层级顺序拼接
            breadcrumb = ' > '.join(level_stack[k] for k in sorted(level_stack))
            current = {
                'level': level,
                'heading': title,
                'breadcrumb': breadcrumb,
                'lines': [line],
            }
        else:
            if current is None:
                # 文档开头的前言内容
                current = {
                    'level': -1,
                    'heading': '前言',
                    'breadcrumb': '前言',
                    'lines': [],
                }
            current['lines'].append(line)

    if current is not None:
        sections.append(current)

    # ---- 第二步：对每节生成块，超长节二次切分 ----
    raw_chunks: List[Dict] = []
    for sec in sections:
        content = '\n'.join(sec['lines']).strip()
        if not content:
            continue
        base_meta = {
            'source': source,
            'breadcrumb': sec['breadcrumb'],
            'heading': sec['heading'],
            'level': sec['level'],
        }
        if len(content) <= chunk_size:
            raw_chunks.append({'text': content, 'metadata': base_meta})
        else:
            sub_chunks = _split_long_section(content, chunk_size, overlap)
            n = len(sub_chunks)
            for idx, sub in enumerate(sub_chunks):
                meta = dict(base_meta)
                meta['chunk_index'] = idx
                meta['total_chunks'] = n
                raw_chunks.append({'text': sub, 'metadata': meta})

    # ---- 第三步：合并短块 ----
    final_chunks = _merge_short_chunks(raw_chunks, min_chunk_size)

    # ---- 第四步：添加全局序号 ----
    for i, chunk in enumerate(final_chunks):
        chunk['chunk_id'] = f"{source}_{i:04d}" if source else f"{i:04d}"

    return final_chunks


# ============================================================
# 文件级处理 & 批量处理
# ============================================================

def process_file(file_path: str, source_name: str = None, **kwargs) -> List[Dict]:
    """提取单个文件文本并分块。"""
    if source_name is None:
        source_name = Path(file_path).stem
    logger.info(f"处理: {file_path}")
    try:
        text = extract_text(file_path)
    except Exception as e:
        logger.error(f"提取失败 {file_path}: {e}")
        return []
    if not text.strip():
        logger.warning(f"文本为空，跳过: {file_path}")
        return []
    chunks = chunk_textbook(text, source=source_name, **kwargs)
    for chunk in chunks:
        chunk['file_path'] = file_path
    return chunks


def batch_process(input_dir: str, output_dir: str, **kwargs):
    """
    批量处理目录下所有支持格式的文档，输出 JSON 分块结果。
    支持格式：.pdf .docx .pptx .txt .md
    """
    input_path = Path(input_dir)
    output_path = Path(output_dir)
    output_path.mkdir(parents=True, exist_ok=True)

    supported = {'.pdf', '.docx', '.pptx', '.txt', '.md'}
    files = [f for f in input_path.glob('*') if f.suffix.lower() in supported]
    if not files:
        logger.warning(f"未找到支持的文档: {input_dir}")
        return

    all_chunks = []
    for file in tqdm(files, desc="分块处理"):
        chunks = process_file(str(file), **kwargs)
        all_chunks.extend(chunks)

        out_file = output_path / f"{file.stem}_chunks.json"
        with open(out_file, 'w', encoding='utf-8') as f:
            json.dump(chunks, f, ensure_ascii=False, indent=2)
        logger.info(f"已保存 {len(chunks)} 个块 -> {out_file}")

    all_out = output_path / "all_chunks.json"
    with open(all_out, 'w', encoding='utf-8') as f:
        json.dump(all_chunks, f, ensure_ascii=False, indent=2)
    logger.info(f"全部 {len(all_chunks)} 个块已保存 -> {all_out}")


# ============================================================
# 入口
# ============================================================

if __name__ == "__main__":
    # 示例：处理单个文件
    sample = r"C:\Users\86182\Downloads\documents"
    output = "./chunks_output"

    batch_process(
        input_dir=sample,
        output_dir=output,
        chunk_size=800,
        overlap=150,
        min_chunk_size=100,
    )
        